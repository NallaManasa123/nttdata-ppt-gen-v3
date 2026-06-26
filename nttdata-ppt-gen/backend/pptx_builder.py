"""
NTT DATA PPTX Builder v4 — Main entry point.
Orchestrates: primitives → renderers → slides → final .pptx bytes.
Template-faithful design matching NTT DATA PowerPoint template.
"""
import io
from pptx import Presentation
from backend.pptx_slides import LAYOUT_MAP, build_cover, build_closing
from backend.pptx_primitives import notes, SW, SH


def _resolve_layout(data: dict) -> str:
    """Determine the correct layout key for a slide."""
    stype  = data.get("slide_type", "content")
    layout = data.get("layout", "two_column")
    title  = (data.get("title") or "").lower()

    if stype == "cover":
        return "cover"
    if stype == "closing":
        return "closing"
    if stype == "agenda" or layout == "agenda" or "agenda" in title:
        return "agenda"
    if layout == "section_divider" or stype == "section_divider":
        return "section_divider"

    # Map LLM layout names → our builder keys
    MAP = {
        "cover":            "cover",
        "agenda":           "agenda",
        "two_column":       "two_column",
        "full_visual":      "full_visual",
        "bullet_list":      "bullet_list",
        "comparison_table": "comparison_table",
        "timeline_visual":  "timeline_visual",
        "process_flow":     "process_flow",
        "kpi_dashboard":    "kpi_dashboard",
        "section_divider":  "section_divider",
        "closing":          "closing",
    }
    return MAP.get(layout, "two_column")


def build_pptx(plan: dict, theme: str = "light") -> bytes:
    """
    Build and return .pptx bytes from a slide plan dict.
    plan must contain: {"slides": [...], "theme": "light"|"dark"}
    """
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    blank_layout = prs.slide_layouts[6]  # blank

    slides_data = plan.get("slides", [])
    total = len(slides_data)

    # Ensure presentation always starts with a cover
    has_cover = any(
        d.get("slide_type") == "cover" or d.get("layout") == "cover"
        for d in slides_data
    )
    if not has_cover and slides_data:
        slides_data = [{"slide_type": "cover", "layout": "cover",
                        "title": plan.get("presentation_title", "NTT DATA"),
                        "subtitle": "Presentation",
                        "bullets": [], "visual": {}, "section": None,
                        "speaker_notes": None}] + slides_data
        total = len(slides_data)

    # Ensure presentation ends with a closing slide
    has_closing = any(
        d.get("slide_type") == "closing" or d.get("layout") == "closing"
        for d in slides_data
    )
    if not has_closing and slides_data:
        slides_data = slides_data + [
            {"slide_type": "closing", "layout": "closing",
             "title": "Thank You",
             "subtitle": "Questions & Discussion",
             "bullets": [], "visual": {}, "section": None,
             "speaker_notes": None}
        ]
        total = len(slides_data)

    for idx, data in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        num   = idx + 1
        resolved = _resolve_layout(data)
        fn = LAYOUT_MAP.get(resolved, LAYOUT_MAP["two_column"])

        try:
            if resolved in ("cover", "closing"):
                fn(slide, data, theme)
            else:
                fn(slide, data, theme, num, total)
        except Exception as e:
            # Fallback: plain bullet list
            from backend.pptx_slides import build_bullet_list
            try:
                build_bullet_list(slide, data, theme, num, total)
            except Exception:
                pass

        notes(slide, data.get("speaker_notes"))

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()
