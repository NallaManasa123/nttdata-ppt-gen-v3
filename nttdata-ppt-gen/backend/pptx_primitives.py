"""NTT DATA PPTX — Primitives, brand colours, template positions, shared helpers."""
import io, base64
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

try:
    from backend.logo_b64 import LOGO_BLUE_B64, LOGO_WHITE_B64
    _LB = base64.b64decode(LOGO_BLUE_B64)
    _LW = base64.b64decode(LOGO_WHITE_B64)
except Exception:
    _LB = _LW = None


# ── Brand palette ───────────────────────────────────────────────────────────
C_NAVY   = RGBColor(0x00,0x18,0x6A)
C_BLUE   = RGBColor(0x00,0x30,0x87)
C_DKBG   = RGBColor(0x00,0x10,0x3A)
C_MID    = RGBColor(0x00,0x50,0xB3)
C_ACCENT = RGBColor(0x00,0xB4,0xD8)
C_WHITE  = RGBColor(0xFF,0xFF,0xFF)
C_LGRAY  = RGBColor(0xF5,0xF7,0xFA)
C_MGRAY  = RGBColor(0xCC,0xD8,0xEA)
C_MUTED  = RGBColor(0x5A,0x6A,0x8A)
C_TEXT   = RGBColor(0x0D,0x1B,0x3E)
C_TEAL   = RGBColor(0x00,0x7A,0x8A)
C_GREEN  = RGBColor(0x0A,0x7C,0x3E)
C_ORANGE = RGBColor(0xE8,0x6A,0x10)
C_RED    = RGBColor(0xC0,0x20,0x20)
C_PURPLE = RGBColor(0x5A,0x20,0x8A)
C_DKTEXT = RGBColor(0xBB,0xCC,0xEE)  # body text on dark bg

COLORS = {
    "blue":C_BLUE,"navy":C_NAVY,"mid":C_MID,"dark":C_DKBG,
    "teal":C_TEAL,"accent":C_ACCENT,"green":C_GREEN,
    "gray":C_MUTED,"orange":C_ORANGE,"red":C_RED,"purple":C_PURPLE,
}

# ── Slide dimensions (16:9 widescreen) ───────────────────────────────────────
SW = Inches(13.333)
SH = Inches(7.5)

# ── Content area (below title, above footer) ──────────────────────────────────
MARGIN    = Inches(0.406)
CONTENT_L = MARGIN
CONTENT_W = SW - MARGIN * 2
CONTENT_T = Inches(1.32)
CONTENT_H = Inches(5.72)

# ── Title area (exact from NTT template) ──────────────────────────────────────
TITLE_L = Inches(0.406)
TITLE_T = Inches(0.365)
TITLE_W = Inches(12.521)
TITLE_H = Inches(0.866)

# ── Cover slide positions (exact from NTT template) ───────────────────────────
CVR_LOGO_L = Inches(0.406)   # logo TOP-LEFT on cover
CVR_LOGO_T = Inches(0.350)
CVR_LOGO_W = Inches(2.143)
CVR_LOGO_H = Inches(0.438)
CVR_TTL_L  = Inches(0.416)
CVR_TTL_T  = Inches(2.667)
CVR_TTL_W  = Inches(5.243)
CVR_TTL_H  = Inches(1.500)
CVR_SUB_L  = Inches(0.416)
CVR_SUB_T  = Inches(4.167)
CVR_SUB_W  = Inches(5.243)
CVR_SUB_H  = Inches(1.060)
CVR_META_L = Inches(0.406)   # presenter/date top-left
CVR_META_T = Inches(0.900)   # below logo on cover

# ── Footer area (exact from NTT template) ─────────────────────────────────────
FOOTER_T        = Inches(7.178)
FOOTER_H        = Inches(0.28)
FOOTER_CR_L     = Inches(0.406)   # copyright LEFT
FOOTER_CR_W     = Inches(5.306)
FOOTER_NUM_L    = Inches(6.500)   # slide number CENTRE
FOOTER_NUM_W    = Inches(0.333)
FOOTER_LOGO_L   = Inches(11.517)  # logo BOTTOM-RIGHT
FOOTER_LOGO_T   = Inches(7.054)
FOOTER_LOGO_W   = Inches(1.410)
FOOTER_LOGO_H   = Inches(0.287)


# ── Logo storage ──────────────────────────────────────────────────────────────
# _LW = white logo → used on dark backgrounds
# _LB = blue logo  → used on light backgrounds
_LOGO_DARK  = _LW   # white logo for dark slides
_LOGO_LIGHT = _LB   # blue logo for light slides

def load_logos(dark_path: str, light_path: str):
    global _LOGO_DARK, _LOGO_LIGHT
    try:
        with open(dark_path,  "rb") as f: _LOGO_DARK  = f.read()
        with open(light_path, "rb") as f: _LOGO_LIGHT = f.read()
    except Exception:
        pass

# ── Colour helper ──────────────────────────────────────────────────────────────
def col(name: str) -> RGBColor:
    return COLORS.get(str(name).lower(), C_MID)

# ── Shape primitives ───────────────────────────────────────────────────────────
def rect(sl, l, t, w, h, fill, lc=None, lpt=0):
    s = sl.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if lc and lpt:
        s.line.color.rgb = lc
        s.line.width = Pt(lpt)
    else:
        s.line.fill.background()
    return s

def rrect(sl, l, t, w, h, fill, lc=None, lpt=0, rad=0.12):
    s = sl.shapes.add_shape(5, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if lc and lpt:
        s.line.color.rgb = lc
        s.line.width = Pt(lpt)
    else:
        s.line.fill.background()
    try:
        if len(s.adjustments) > 0:
            s.adjustments[0] = rad
    except Exception:
        pass
    return s

def oval(sl, l, t, w, h, fill=None, lc=None, lpt=0):
    s = sl.shapes.add_shape(9, l, t, w, h)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if lc and lpt:
        s.line.color.rgb = lc
        s.line.width = Pt(lpt)
    else:
        s.line.fill.background()
    return s

def arrow(sl, x1, y1, x2, y2, color=None, w=1.5):
    c = sl.shapes.add_connector(2, x1, y1, x2, y2)
    c.line.color.rgb = color or C_ACCENT
    c.line.width = Pt(w)
    return c

def txt(sl, text, l, t, w, h,
        font="Calibri", size=12, bold=False,
        color=None, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    if color is None:
        color = C_TEXT
    tb = sl.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text = text
    r.font.name   = font
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb

def bullets(sl, items, l, t, w, h, theme, size=13):
    if not items: return
    tc = C_TEXT if theme == "light" else C_WHITE
    tb = sl.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for b in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(6)
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = "\u25b8  " + b
        r.font.name  = "Calibri"
        r.font.size  = Pt(size)
        r.font.color.rgb = tc

def logo(sl, theme, l, t, w, h):
    data = _LOGO_DARK if theme == "dark" else _LOGO_LIGHT
    if data is None: return
    try:
        sl.shapes.add_picture(io.BytesIO(data), l, t, width=w, height=h)
    except Exception:
        pass

def footer(sl, theme, num=None, total=None):
    """Footer: copyright LEFT, slide number CENTRE, logo RIGHT."""
    cc = C_MUTED if theme == "light" else RGBColor(0x88, 0x9A, 0xB8)
    txt(sl, "\u00a9 2026 NTT DATA Business Solutions",
        FOOTER_CR_L, FOOTER_T, FOOTER_CR_W, FOOTER_H,
        size=8, color=cc)
    if num and total:
        nc = C_MUTED if theme == "light" else RGBColor(0x77, 0x88, 0xAA)
        txt(sl, f"{num} / {total}",
            FOOTER_NUM_L, FOOTER_T, FOOTER_NUM_W, FOOTER_H,
            size=8, color=nc, align=PP_ALIGN.CENTER)
    logo(sl, theme, FOOTER_LOGO_L, FOOTER_LOGO_T, FOOTER_LOGO_W, FOOTER_LOGO_H)

def bg(sl, theme):
    fill = C_LGRAY if theme == "light" else RGBColor(0x04, 0x10, 0x30)
    rect(sl, 0, 0, SW, SH, fill)

def dark_bg(sl):
    rect(sl, 0, 0, SW, SH, C_DKBG)


# ── Content title area (short alias for TITLE_* constants) ───────────────────
CT_TTL_L = TITLE_L          # Inches(0.406)
CT_TTL_T = TITLE_T          # Inches(0.365)
CT_TTL_W = TITLE_W          # Inches(12.521)
CT_TTL_H = TITLE_H          # Inches(0.866)

# ── Content area aliases ──────────────────────────────────────────────────────
CA_L = CONTENT_L            # Inches(0.406)
CA_T = CONTENT_T            # Inches(1.32)
CA_W = CONTENT_W            # Inches(12.521)
CA_H = CONTENT_H            # Inches(5.72)

# ── Footer logo aliases (short names used by pptx_slides.py) ─────────────────
FT_LOGO_L = FOOTER_LOGO_L   # Inches(11.517)
FT_LOGO_T = FOOTER_LOGO_T   # Inches(7.054)
FT_LOGO_W = FOOTER_LOGO_W   # Inches(1.410)
FT_LOGO_H = FOOTER_LOGO_H   # Inches(0.287)

# ── Closing slide logo (centred horizontally) ─────────────────────────────────
CL_LOGO_W = Inches(2.143)
CL_LOGO_H = Inches(0.438)
CL_LOGO_L = (SW - CL_LOGO_W) / 2   # horizontally centred
CL_LOGO_T = Inches(4.6)


# ── Slide title helper ────────────────────────────────────────────────────────
def slide_title(sl, title, theme="light"):
    """Render a slide title in the standard NTT DATA title area."""
    tc = C_TEXT if theme == "light" else C_WHITE
    txt(sl, title, TITLE_L, TITLE_T, TITLE_W, TITLE_H,
        font="Calibri", size=28, bold=True, color=tc)


# ── Speaker notes helper ──────────────────────────────────────────────────────
def notes(sl, text):
    """Write speaker notes to a slide's notes pane."""
    if not text:
        return
    try:
        notes_slide = sl.notes_slide
        notes_slide.notes_text_frame.text = str(text)
    except Exception:
        pass


# ── Backward-compat aliases (used by pptx_renderers.py) ──────────────────────
_rect    = rect
_rrect   = rrect
_oval    = oval
_arrow   = arrow
_txt     = txt
_bullets = bullets
_col     = col
_logo    = logo
